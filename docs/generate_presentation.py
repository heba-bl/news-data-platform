"""
THEME CLAIR STYLE AI — News Data Platform
22 slides : fond blanc, bleu AI, typographie moderne
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import sys
sys.stdout.reconfigure(encoding='utf-8')

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── PALETTE CLAIRE STYLE AI ─────────────────────────────────
BG        = RGBColor(0xF8, 0xF9, 0xFF)   # blanc légèrement bleuté
BG2       = RGBColor(0xEE, 0xF2, 0xFF)   # bleu très clair
CARD      = RGBColor(0xFF, 0xFF, 0xFF)   # blanc pur
BORDER    = RGBColor(0xE2, 0xE8, 0xF0)   # gris-bleu clair
C_INK     = RGBColor(0x0F, 0x17, 0x2A)   # quasi-noir
C_DARK    = RGBColor(0x1E, 0x29, 0x4D)   # bleu foncé
C_TEXT    = RGBColor(0x33, 0x44, 0x66)   # texte corps
C_MUTED   = RGBColor(0x7A, 0x8A, 0xAA)   # gris discret
C_BLUE    = RGBColor(0x24, 0x63, 0xEB)   # bleu AI principal
C_INDIGO  = RGBColor(0x4F, 0x46, 0xE5)   # indigo
C_TEAL    = RGBColor(0x06, 0xB6, 0xD4)   # cyan
C_GREEN   = RGBColor(0x10, 0xB9, 0x81)   # vert
C_ORANGE  = RGBColor(0xF5, 0x9E, 0x0B)   # ambre
C_RED     = RGBColor(0xEF, 0x44, 0x44)   # rouge
C_PURPLE  = RGBColor(0x8B, 0x5C, 0xF6)   # violet
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

blank = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(blank)


def bg(slide, color=BG):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def rect(slide, l, t, w, h, color, radius=False):
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def txt(slide, text, l, t, w, h,
        size=16, bold=False, color=C_TEXT,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


# Barre latérale gauche colorée (signature AI)
def left_bar(slide, color=C_BLUE):
    rect(slide, 0, 0, 0.18, 7.5, color)


# En-tête standard slide
def header(slide, title, sub=None, accent=C_BLUE):
    rect(slide, 0.18, 0, 13.15, 1.25, BG2)
    rect(slide, 0.18, 1.25, 13.15, 0.055, accent)
    txt(slide, title, 0.45, 0.1, 12.5, 0.75, size=27, bold=True, color=C_DARK)
    if sub:
        txt(slide, sub, 0.45, 0.8, 12.5, 0.42, size=13, color=C_MUTED, italic=True)


def page_num(slide, n, total=22):
    txt(slide, f"{n} / {total}", 12.4, 7.1, 0.9, 0.35, size=10,
        color=C_MUTED, align=PP_ALIGN.RIGHT)


def chip(slide, label, l, t, w, h, bg_color, fg_color=C_WHITE):
    rect(slide, l, t, w, h, bg_color)
    txt(slide, label, l + 0.08, t + 0.04, w - 0.16, h - 0.08,
        size=11, bold=True, color=fg_color, align=PP_ALIGN.CENTER)


def card(slide, title, body, l, t, w, h, accent=C_BLUE, title_size=14, body_size=12):
    rect(slide, l, t, w, h, CARD)
    rect(slide, l, t, 0.055, h, accent)
    rect(slide, l, t, w, h, BORDER)          # fine bordure
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.background()
    shp.line.color.rgb = BORDER
    shp.line.width = Pt(0.75)
    txt(slide, title, l + 0.18, t + 0.1, w - 0.28, 0.4,
        size=title_size, bold=True, color=accent)
    txt(slide, body,  l + 0.18, t + 0.55, w - 0.28, h - 0.65,
        size=body_size, color=C_TEXT)


# ════════════════════════════════════════════════════════════
# SLIDE 1 — TITRE
# ════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 7.5, C_DARK)
rect(s, 0, 0, 0.22, 7.5, C_BLUE)
rect(s, 0.22, 5.7, 13.11, 1.8, RGBColor(0x14, 0x1E, 0x3A))

# Logo-style "AI" badge
rect(s, 0.55, 0.55, 1.3, 0.55, C_BLUE)
txt(s, "AI  PROJECT", 0.6, 0.58, 1.25, 0.45, size=11, bold=True,
    color=C_WHITE, align=PP_ALIGN.CENTER)

txt(s, "ARCHITECTURE DE DONNÉES", 0.55, 1.4, 12.2, 0.75,
    size=14, bold=False, color=C_TEAL, align=PP_ALIGN.LEFT)
txt(s, "Big Data pour l'Analyse", 0.55, 2.05, 12.2, 1.0,
    size=42, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
txt(s, "des Articles de Presse", 0.55, 2.95, 12.2, 0.85,
    size=42, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

txt(s, "Kafka  ·  Airflow  ·  MinIO  ·  PostgreSQL  ·  Grafana  ·  Docker",
    0.55, 4.1, 12.2, 0.45, size=13, color=C_MUTED, align=PP_ALIGN.LEFT)

for i, (lbl, w_) in enumerate([("Nom :", 1.5), ("Groupe :", 1.5), ("Date :", 1.5)]):
    x = 0.55 + i * 4.15
    txt(s, lbl,          x,       5.9, 1.1, 0.4, size=12, bold=True,  color=C_TEAL)
    txt(s, "_______________", x+1.0, 5.9, 2.8, 0.4, size=12, color=C_MUTED)

page_num(s, 1)


# ════════════════════════════════════════════════════════════
# SLIDE 2 — SOMMAIRE
# ════════════════════════════════════════════════════════════
s = add_slide(); bg(s); left_bar(s, C_INDIGO)
header(s, "Sommaire", "Plan de la présentation", C_INDIGO)

items = [
    ("01", "Contexte & Objectifs",      "Problématique · Buts du projet"),
    ("02", "Sources de Données",        "Sites marocains et internationaux"),
    ("03", "Architecture Globale",      "Pipeline bout en bout"),
    ("04", "Ingestion des Données",     "Kafka Streaming + Airflow Batch"),
    ("05", "Data Lake Médaillon",       "Bronze → Silver → Gold (MinIO)"),
    ("06", "Transformation ETL",        "Nettoyage · Normalisation · Enrichissement"),
    ("07", "Orchestration Airflow",     "3 DAGs planifiés"),
    ("08", "Data Warehouse",            "Modèle étoile PostgreSQL"),
    ("09", "Visualisation Grafana",     "13 panels analytiques"),
    ("10", "Qualité & Gouvernance",     "Contrôles · Lineage · Traçabilité"),
    ("11", "Déploiement & Monitoring",  "Docker · Kubernetes · Prometheus"),
    ("12", "Résultats & Livrables",     "Démonstration · Conclusion"),
]

cols = [items[:6], items[6:]]
accents = [C_BLUE, C_INDIGO]
for ci, col_items in enumerate(cols):
    x = 0.4 + ci * 6.5
    for i, (num, title, sub) in enumerate(col_items):
        y = 1.5 + i * 0.88
        rect(s, x, y, 6.2, 0.76, CARD)
        shp = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(6.2), Inches(0.76))
        shp.fill.background(); shp.line.color.rgb = BORDER; shp.line.width = Pt(0.6)
        chip(s, num, x + 0.07, y + 0.13, 0.5, 0.48,
             accents[ci] if i % 2 == 0 else C_TEAL)
        txt(s, title, x + 0.7, y + 0.07, 5.3, 0.35, size=13, bold=True, color=C_DARK)
        txt(s, sub,   x + 0.7, y + 0.4,  5.3, 0.3,  size=10, color=C_MUTED, italic=True)

page_num(s, 2)


# ════════════════════════════════════════════════════════════
# SLIDE 3 — CONTEXTE
# ════════════════════════════════════════════════════════════
s = add_slide(); bg(s); left_bar(s, C_TEAL)
header(s, "Contexte", "Pourquoi une plateforme d'analyse de presse ?", C_TEAL)

points = [
    ("📰", "Volume massif",    C_BLUE,   "Des milliers d'articles publiés chaque jour,\nimpossibles à analyser manuellement."),
    ("📉", "Tendances floues", C_INDIGO, "Sans outil, identifier les sujets dominants\nprend des heures de lecture."),
    ("⚡", "Besoin temps réel",C_TEAL,   "Les événements évoluent vite — une\nanalyse tardive perd toute valeur."),
    ("🌍", "Sources multiples",C_PURPLE, "Agréger presse marocaine et\ninternationale nécessite une infra unifiée."),
]

for i, (icon, title, accent, body) in enumerate(points):
    col = i % 2; row = i // 2
    x = 0.4 + col * 6.5; y = 1.55 + row * 2.6
    rect(s, x, y, 6.2, 2.35, CARD)
    shp = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(6.2), Inches(2.35))
    shp.fill.background(); shp.line.color.rgb = accent; shp.line.width = Pt(1.5)
    rect(s, x, y, 6.2, 0.055, accent)
    txt(s, icon,  x+0.2,  y+0.2,  0.6,  0.6, size=24, align=PP_ALIGN.CENTER)
    txt(s, title, x+0.9,  y+0.2,  5.1,  0.5, size=16, bold=True, color=accent)
    txt(s, body,  x+0.2,  y+0.85, 5.8,  1.3, size=13, color=C_TEXT)

page_num(s, 3)


# ════════════════════════════════════════════════════════════
# SLIDE 4 — OBJECTIFS
# ════════════════════════════════════════════════════════════
s = add_slide(); bg(s); left_bar(s, C_GREEN)
header(s, "Objectifs du Projet", "6 axes techniques couverts", C_GREEN)

objs = [
    ("01", C_BLUE,   "Collecter automatiquement",  "Scraping 5 sources toutes les heures"),
    ("02", C_TEAL,   "Pipeline Big Data",           "Kafka streaming + Airflow orchestration"),
    ("03", C_GREEN,  "Architecture Médaillon",      "Bronze → Silver → Gold dans MinIO"),
    ("04", C_ORANGE, "Data Warehouse",              "Modèle dimensionnel PostgreSQL"),
    ("05", C_RED,    "Qualité & Gouvernance",       "Contrôles automatisés + data lineage"),
    ("06", C_PURPLE, "Visualisation",               "Grafana 13 panels + Metabase BI"),
]

for i, (num, color, title, body) in enumerate(objs):
    col = i % 3; row = i // 3
    x = 0.4 + col * 4.3; y = 1.55 + row * 2.7
    rect(s, x, y, 4.1, 2.45, CARD)
    shp = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(4.1), Inches(2.45))
    shp.fill.background(); shp.line.color.rgb = BORDER; shp.line.width = Pt(0.75)
    rect(s, x, y, 4.1, 0.055, color)
    rect(s, x+0.15, y+0.18, 0.6, 0.6, color)
    txt(s, num,   x+0.22,  y+0.22, 0.48, 0.48, size=17, bold=True,
        color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, x+0.9,   y+0.18, 3.05, 0.45, size=14, bold=True, color=color)
    txt(s, body,  x+0.18,  y+0.92, 3.78, 1.35, size=12, color=C_TEXT)

page_num(s, 4)


# ════════════════════════════════════════════════════════════
# SLIDE 5 — SOURCES DE DONNÉES
# ════════════════════════════════════════════════════════════
s = add_slide(); bg(s); left_bar(s, C_ORANGE)
header(s, "Sources de Données", "5 sites — marocains & internationaux", C_ORANGE)

src_ma  = [("🇲🇦  Hespress",  "ar / fr", "Politique · Sport · Économie"),
           ("🇲🇦  Akhbarona", "ar",      "Actualité générale")]
src_int = [("🇬🇧  BBC News",  "en", "Monde · Tech · Science · Sport"),
           ("🇺🇸  CNN",       "en", "Monde · Politique · Business"),
           ("🇬🇧  Reuters",   "en", "Finance · Politique · Marchés")]

rect(s, 0.4,  1.5, 0.95, 0.45, BG2)
txt(s, "Marocains", 0.45, 1.52, 0.88, 0.4, size=12, bold=True, color=C_ORANGE)
for i, (name, lang, cats) in enumerate(src_ma):
    y = 2.05 + i * 1.45
    rect(s, 0.4, y, 5.8, 1.28, CARD)
    shp = s.shapes.add_shape(1, Inches(0.4), Inches(y), Inches(5.8), Inches(1.28))
    shp.fill.background(); shp.line.color.rgb = C_ORANGE; shp.line.width = Pt(1.2)
    txt(s, name,  0.6, y+0.1,  5.4, 0.45, size=15, bold=True, color=C_DARK)
    chip(s, lang, 0.6, y+0.65, 0.7, 0.32, C_ORANGE)
    txt(s, cats,  1.4, y+0.67, 4.6, 0.35, size=11, color=C_MUTED)

rect(s, 6.8,  1.5, 1.15, 0.45, BG2)
txt(s, "Internationaux", 6.85, 1.52, 1.1, 0.4, size=12, bold=True, color=C_BLUE)
for i, (name, lang, cats) in enumerate(src_int):
    y = 2.05 + i * 1.45
    rect(s, 6.8, y, 6.1, 1.28, CARD)
    shp = s.shapes.add_shape(1, Inches(6.8), Inches(y), Inches(6.1), Inches(1.28))
    shp.fill.background(); shp.line.color.rgb = C_BLUE; shp.line.width = Pt(1.2)
    txt(s, name,  7.0, y+0.1,  5.7, 0.45, size=15, bold=True, color=C_DARK)
    chip(s, lang, 7.0, y+0.65, 0.55, 0.32, C_BLUE)
    txt(s, cats,  7.65, y+0.67, 5.0, 0.35, size=11, color=C_MUTED)

rect(s, 0.4, 6.08, 12.5, 1.08, BG2)
fields = "Données collectées :   Titre  ·  Auteur  ·  Date  ·  Catégorie  ·  Contenu  ·  Source  ·  URL  ·  Langue détectée"
txt(s, fields, 0.55, 6.16, 12.2, 0.5, size=13, bold=True,
    color=C_BLUE, align=PP_ALIGN.CENTER)

page_num(s, 5)


# ════════════════════════════════════════════════════════════
# SLIDE 6 — ARCHITECTURE GLOBALE
# ════════════════════════════════════════════════════════════
s = add_slide(); bg(s); left_bar(s, C_BLUE)
header(s, "Architecture Globale", "Pipeline de données de bout en bout", C_BLUE)

steps  = ["SCRAPING\nPython\nBeautifulSoup",
          "KAFKA\nStreaming\nnews-articles",
          "MINIO\nBronze→Silver\n→Gold",
          "ETL\nPython\nPandas",
          "POSTGRESQL\nData\nWarehouse",
          "GRAFANA\n13 Panels\nLive"]
colors = [C_TEAL, C_BLUE, C_ORANGE, C_INDIGO, C_RED, C_GREEN]
icons  = ["🕷", "📨", "🗄", "⚙️", "🏛", "📊"]

box_w = 1.75; gap = 0.28
total = len(steps) * box_w + (len(steps)-1) * gap
sx = (13.33 - total) / 2

for i, (step, color, icon) in enumerate(zip(steps, colors, icons)):
    x = sx + i * (box_w + gap)
    rect(s, x, 2.1, box_w, 2.1, CARD)
    shp = s.shapes.add_shape(1, Inches(x), Inches(2.1), Inches(box_w), Inches(2.1))
    shp.fill.background(); shp.line.color.rgb = color; shp.line.width = Pt(1.8)
    rect(s, x, 2.1, box_w, 0.065, color)
    txt(s, icon, x+0.5, 2.2, 0.75, 0.6, size=22, align=PP_ALIGN.CENTER)
    txt(s, step, x+0.05, 2.85, box_w-0.1, 1.2, size=11, bold=True,
        color=color, align=PP_ALIGN.CENTER)
    if i < len(steps)-1:
        txt(s, "→", x+box_w+0.03, 2.85, 0.22, 0.6, size=18,
            bold=True, color=C_MUTED, align=PP_ALIGN.CENTER)

rect(s, 0.4, 4.55, 12.5, 0.55, BG2)
tech = "Docker Compose  ·  11 services  ·  Apache Airflow  ·  Apache Kafka 7.4  ·  Python 3.11  ·  PostgreSQL 15  ·  Grafana 13"
txt(s, tech, 0.55, 4.62, 12.2, 0.42, size=12, color=C_MUTED,
    italic=True, align=PP_ALIGN.CENTER)

# KPIs réels
kpis = [("900+", "Articles Bronze"), ("173", "Chargés DWH"),
        ("5", "Sources"), ("131", "Mots-clés"), ("3", "Pays"), ("75.7%", "Qualité")]
for i, (val, lbl) in enumerate(kpis):
    x = 0.4 + i * 2.09
    rect(s, x, 5.3, 1.97, 1.9, CARD)
    shp = s.shapes.add_shape(1, Inches(x), Inches(5.3), Inches(1.97), Inches(1.9))
    shp.fill.background(); shp.line.color.rgb = BORDER; shp.line.width = Pt(0.6)
    txt(s, val, x+0.05, 5.42, 1.88, 0.72, size=22, bold=True,
        color=C_BLUE, align=PP_ALIGN.CENTER)
    txt(s, lbl, x+0.05, 6.15, 1.88, 0.45, size=10, color=C_MUTED,
        align=PP_ALIGN.CENTER)

page_num(s, 6)


# ════════════════════════════════════════════════════════════
# SLIDE 7 — DATA INGESTION
# ════════════════════════════════════════════════════════════
s = add_slide(); bg(s); left_bar(s, C_TEAL)
header(s, "Data Ingestion", "Deux modes d'ingestion complémentaires", C_TEAL)

for ci, (title, color, pts, detail) in enumerate([
    ("⏱  Batch Ingestion", C_BLUE,
     ["Scraping planifié toutes les heures",
      "Orchestré par Apache Airflow DAG",
      "DAG : scraping_batch (@hourly)",
      "5 sources scrapées en parallèle",
      "Chaque article envoyé vers Kafka"],
     "Apache Airflow\nscraping_batch"),
    ("⚡  Streaming Ingestion", C_TEAL,
     ["Chaque article = 1 message Kafka",
      "Topic : news-articles (3 partitions)",
      "Consumer lit en continu",
      "Stockage immédiat Bronze (MinIO)",
      "Rétention : 7 jours"],
     "Apache Kafka\nnews-articles"),
]):
    x = 0.4 + ci * 6.5
    rect(s, x, 1.5, 6.2, 5.65, CARD)
    shp = s.shapes.add_shape(1, Inches(x), Inches(1.5), Inches(6.2), Inches(5.65))
    shp.fill.background(); shp.line.color.rgb = color; shp.line.width = Pt(1.5)
    rect(s, x, 1.5, 6.2, 0.065, color)
    txt(s, title, x+0.2, 1.6, 5.8, 0.52, size=18, bold=True, color=color)
    for i, pt in enumerate(pts):
        txt(s, f"✓  {pt}", x+0.2, 2.3+i*0.7, 5.8, 0.6, size=13, color=C_TEXT)
    rect(s, x+0.2, 5.6, 5.8, 1.25, BG2)
    txt(s, "Technologie", x+0.35, 5.68, 5.5, 0.35, size=10, bold=True, color=C_MUTED)
    txt(s, detail,       x+0.35, 6.0,  5.5, 0.6,  size=14, bold=True, color=color)

page_num(s, 7)


# ════════════════════════════════════════════════════════════
# SLIDE 8 — DATA LAKE
# ════════════════════════════════════════════════════════════
s = add_slide(); bg(s); left_bar(s, C_ORANGE)
header(s, "Data Lake — MinIO", "Object Storage open-source compatible S3", C_ORANGE)

rect(s, 0.4, 1.5, 12.5, 0.95, BG2)
txt(s, "MinIO — Déployé via Docker · API S3 · Interface console web · Volumes persistants · 3 buckets",
    0.55, 1.62, 12.2, 0.42, size=13, color=C_TEXT, align=PP_ALIGN.CENTER)

layers = [
    ("BRONZE",  C_RED,    "bronze/raw/",     "Format : JSON brut",
     ["Articles tels que scrapés", "Historique complet", "Aucune transformation", "900+ fichiers"]),
    ("SILVER",  C_TEAL,   "silver/cleaned/", "Format : Parquet (Snappy)",
     ["HTML supprimé", "Texte normalisé", "Langue détectée", "Dédupliqué par URL"]),
    ("GOLD",    C_ORANGE, "gold/",           "Format : Parquet optimisé",
     ["Agrégations pré-calculées", "Articles/source/jour/pays", "Top mots-clés", "20 fichiers"]),
]

for i, (name, color, path, fmt, pts) in enumerate(layers):
    x = 0.4 + i * 4.3
    rect(s, x, 2.62, 4.1, 4.55, CARD)
    shp = s.shapes.add_shape(1, Inches(x), Inches(2.62), Inches(4.1), Inches(4.55))
    shp.fill.background(); shp.line.color.rgb = BORDER; shp.line.width = Pt(0.75)
    rect(s, x, 2.62, 4.1, 0.065, color)
    chip(s, name, x+0.15, 2.72, 1.1, 0.42, color)
    txt(s, path, x+0.15, 3.25, 3.8, 0.38, size=12, color=C_BLUE, italic=True)
    txt(s, fmt,  x+0.15, 3.65, 3.8, 0.35, size=11, color=C_MUTED)
    for j, pt in enumerate(pts):
        txt(s, f"▶  {pt}", x+0.15, 4.12+j*0.65, 3.8, 0.58, size=12, color=C_TEXT)

    if i < 2:
        txt(s, "→", x+4.0, 4.55, 0.35, 0.6, size=22, bold=True,
            color=C_MUTED, align=PP_ALIGN.CENTER)

page_num(s, 8)


# ════════════════════════════════════════════════════════════
# SLIDE 9 — ARCHITECTURE MÉDAILLON
# ════════════════════════════════════════════════════════════
s = add_slide(); bg(s); left_bar(s, C_ORANGE)
header(s, "Architecture Médaillon", "Bronze  →  Silver  →  Gold", C_ORANGE)

medals = [
    ("🟫 BRONZE", C_RED,    "Données Brutes",
     ["JSON bruts depuis Kafka Consumer",
      "Historique complet conservé",
      "Aucune modification des données",
      "Partitionné par date YYYY/MM/DD",
      "900+ fichiers — 5 sources"]),
    ("⬜ SILVER",  C_TEAL,  "Données Nettoyées",
     ["Suppression des balises HTML",
      "Normalisation Unicode + espaces",
      "Détection langue (langdetect)",
      "Déduplication par URL unique",
      "13 fichiers Parquet compressés"]),
    ("🥇 GOLD",   C_ORANGE, "Données Analytiques",
     ["Agrégations pré-calculées",
      "Articles par source / pays / jour",
      "Top mots-clés (stopwords retirés)",
      "Distribution des langues",
      "20 fichiers Parquet — prêt DWH"]),
]

for i, (title, color, subtitle, pts) in enumerate(medals):
    x = 0.35 + i * 4.33
    rect(s, x, 1.55, 4.1, 5.6, CARD)
    shp = s.shapes.add_shape(1, Inches(x), Inches(1.55), Inches(4.1), Inches(5.6))
    shp.fill.background(); shp.line.color.rgb = color; shp.line.width = Pt(1.5)
    rect(s, x, 1.55, 4.1, 0.07, color)
    txt(s, title,    x+0.15, 1.65, 3.8, 0.52, size=18, bold=True, color=color)
    txt(s, subtitle, x+0.15, 2.25, 3.8, 0.38, size=12, color=C_MUTED, italic=True)
    for j, pt in enumerate(pts):
        txt(s, f"✓  {pt}", x+0.15, 2.78+j*0.68, 3.8, 0.6, size=12, color=C_TEXT)
    if i < 2:
        txt(s, "⟹", x+4.0, 4.0, 0.35, 0.7, size=24, bold=True,
            color=C_MUTED, align=PP_ALIGN.CENTER)

rect(s, 0.35, 7.12, 12.65, 0.32, BG2)
txt(s, "ETL Python :  bronze_to_silver.py  ·  silver_to_gold.py  |  Orchestration : Apache Airflow",
    0.5, 7.16, 12.3, 0.26, size=11, color=C_BLUE, align=PP_ALIGN.CENTER)

page_num(s, 9)


# ════════════════════════════════════════════════════════════
# SLIDE 10 — TRANSFORMATION ETL
# ════════════════════════════════════════════════════════════
s = add_slide(); bg(s); left_bar(s, C_PURPLE)
header(s, "Transformation des Données", "Pipeline ETL Python — 6 étapes", C_PURPLE)

etl_steps = [
    ("1", C_TEAL,   "Lecture Bronze",     "Lecture fichiers JSON\ndepuis MinIO"),
    ("2", C_GREEN,  "Nettoyage HTML",      "BeautifulSoup retire\ntoutes les balises"),
    ("3", C_BLUE,   "Normalisation",       "Unicode · Espaces\nCaractères spéciaux"),
    ("4", C_PURPLE, "Détection Langue",    "langdetect\n(EN / AR / FR)"),
    ("5", C_ORANGE, "Déduplication",       "Hash URL unique\nanti-doublons"),
    ("6", C_RED,    "Écriture Silver",     "Parquet Snappy\nvers MinIO silver/"),
]

for i, (num, color, title, body) in enumerate(etl_steps):
    col = i % 3; row = i // 3
    x = 0.35 + col * 4.33; y = 1.58 + row * 2.7
    rect(s, x, y, 4.1, 2.48, CARD)
    shp = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(4.1), Inches(2.48))
    shp.fill.background(); shp.line.color.rgb = BORDER; shp.line.width = Pt(0.75)
    rect(s, x, y, 4.1, 0.07, color)
    rect(s, x+0.15, y+0.18, 0.55, 0.55, color)
    txt(s, num,   x+0.21,  y+0.22, 0.44, 0.44, size=16, bold=True,
        color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, x+0.85,  y+0.18, 3.1,  0.45, size=14, bold=True, color=color)
    txt(s, body,  x+0.85,  y+0.7,  3.1,  1.55, size=13, color=C_TEXT)

page_num(s, 10)


# ════════════════════════════════════════════════════════════
# SLIDE 11 — ORCHESTRATION AIRFLOW
# ════════════════════════════════════════════════════════════
s = add_slide(); bg(s); left_bar(s, C_GREEN)
header(s, "Orchestration — Apache Airflow", "3 DAGs planifiés et supervisés", C_GREEN)

dags = [
    ("DAG 1", C_BLUE,   "scraping_batch",          "@hourly",
     ["Toutes les heures automatiquement",
      "Lance 5 scrapers en parallèle",
      "Chaque article → KafkaProducer",
      "Consumer → MinIO Bronze"]),
    ("DAG 2", C_ORANGE, "transformation_pipeline", "@daily",
     ["Lit les fichiers Bronze (MinIO)",
      "ETL : nettoyage + normalisation",
      "Détection langue + dédup",
      "Parquet Silver + Gold → MinIO"]),
    ("DAG 3", C_GREEN,  "warehouse_load",          "@daily",
     ["Lit Silver + Gold depuis MinIO",
      "Charge fact_articles + dimensions",
      "Met à jour tables agrégées",
      "Log qualité → data_quality_log"]),
]

for i, (label, color, dag_id, schedule, pts) in enumerate(dags):
    x = 0.35 + i * 4.33
    rect(s, x, 1.52, 4.1, 5.65, CARD)
    shp = s.shapes.add_shape(1, Inches(x), Inches(1.52), Inches(4.1), Inches(5.65))
    shp.fill.background(); shp.line.color.rgb = BORDER; shp.line.width = Pt(0.75)
    rect(s, x, 1.52, 4.1, 0.07, color)
    chip(s, label, x+0.15, 1.63, 0.85, 0.38, color)
    txt(s, dag_id,   x+0.15, 2.12, 3.8, 0.48, size=15, bold=True, color=C_DARK)
    rect(s, x+0.15, 2.7, 3.8, 0.42, BG2)
    txt(s, f"⏰  Schedule : {schedule}", x+0.25, 2.74, 3.6, 0.34, size=12,
        bold=True, color=color)
    for j, pt in enumerate(pts):
        txt(s, f"▶  {pt}", x+0.15, 3.25+j*0.68, 3.8, 0.6, size=12, color=C_TEXT)

rect(s, 0, 7.08, 13.33, 0.42, BG2)
txt(s, "Interface Airflow  →  http://localhost:8080  |  admin / admin123  |  Airflow 2.8",
    0, 7.13, 13.33, 0.32, size=11, color=C_BLUE, align=PP_ALIGN.CENTER)

page_num(s, 11)


# ════════════════════════════════════════════════════════════
# SLIDE 12 — DATA WAREHOUSE
# ════════════════════════════════════════════════════════════
s = add_slide(); bg(s); left_bar(s, C_RED)
header(s, "Data Warehouse — PostgreSQL", "Modèle en étoile (Star Schema)", C_RED)

# Schéma étoile
rect(s, 0.4, 1.5, 7.0, 5.65, CARD)
shp = s.shapes.add_shape(1, Inches(0.4), Inches(1.5), Inches(7.0), Inches(5.65))
shp.fill.background(); shp.line.color.rgb = BORDER; shp.line.width = Pt(0.75)
txt(s, "Schéma en Étoile", 0.6, 1.6, 6.7, 0.42, size=14, bold=True, color=C_RED)

schema_lines = [
    "            ┌─────────────┐",
    "            │  dim_date   │",
    "            │  date_id PK │",
    "            │  day/month  │",
    "            │  year/qtr   │",
    "            └──────┬──────┘",
    "                   │",
    "┌──────────┐        │         ┌─────────────┐",
    "│dim_source├────fact_articles─┤dim_category │",
    "└──────────┘   title · author └─────────────┘",
    "               content · url",
    "               published_at · language",
    "               word_count · source_id",
]
for i, line in enumerate(schema_lines):
    txt(s, line, 0.55, 2.1+i*0.35, 6.7, 0.35, size=10, color=C_TEXT)

# Tables analytiques
rect(s, 7.8, 1.5, 5.1, 5.65, CARD)
shp = s.shapes.add_shape(1, Inches(7.8), Inches(1.5), Inches(5.1), Inches(5.65))
shp.fill.background(); shp.line.color.rgb = BORDER; shp.line.width = Pt(0.75)
txt(s, "Tables Analytiques", 8.0, 1.6, 4.8, 0.42, size=14, bold=True, color=C_ORANGE)

agg = [
    ("agg_articles_by_source",   C_BLUE,   "Articles par source"),
    ("agg_articles_by_day",      C_TEAL,   "Tendance temporelle"),
    ("agg_articles_by_country",  C_GREEN,  "Articles par pays"),
    ("agg_articles_by_category", C_ORANGE, "Par catégorie / thème"),
    ("agg_top_keywords",         C_PURPLE, "Top 131 mots-clés"),
    ("agg_language_distribution",C_RED,    "Distribution langues"),
    ("pipeline_runs",            C_MUTED,  "Historique pipelines"),
    ("data_quality_log",         C_MUTED,  "Rapport qualité"),
]
for i, (tbl, color, desc) in enumerate(agg):
    y = 2.12 + i * 0.62
    rect(s, 8.0, y, 4.7, 0.55, BG2)
    txt(s, tbl,  8.1,  y+0.07, 4.5, 0.24, size=9,  bold=True, color=color)
    txt(s, desc, 8.1,  y+0.3,  4.5, 0.22, size=9,  color=C_MUTED)

page_num(s, 12)


# ════════════════════════════════════════════════════════════
# SLIDE 13 — VISUALISATION GRAFANA
# ════════════════════════════════════════════════════════════
s = add_slide(); bg(s); left_bar(s, C_TEAL)
header(s, "Visualisation — Grafana", "13 panels analytiques — données réelles", C_TEAL)

panels = [
    ("🟢 Total Articles",       "173",           C_GREEN),
    ("🔵 Sources Actives",      "5",             C_BLUE),
    ("🟠 Mots Clés",            "131",           C_ORANGE),
    ("🟣 Taux Qualité",         "75.7%",         C_PURPLE),
    ("⬜ Dernier Scraping",     "14/05/2026",    C_TEAL),
    ("📰 Par Source",           "BBC·Reuters·CNN",C_BLUE),
    ("🌍 Par Pays",             "UK·USA·Morocco",C_GREEN),
    ("🗣 Langues",              "EN:171 · AR:2", C_ORANGE),
    ("📈 Tendance / Jour",      "21 jours hist.",C_RED),
    ("🔑 Top Mots Clés",        "minister·labour",C_TEAL),
    ("📂 Catégories",           "World·Politics",C_PURPLE),
    ("⚙️  Pipeline Runs",       "4 pipelines",   C_GREEN),
    ("✅ Qualité rapport",      "passed/failed",  C_ORANGE),
]

for i, (name, val, color) in enumerate(panels):
    col = i % 5; row = i // 5
    x = 0.3 + col * 2.55; y = 1.55 + row * 1.92
    rect(s, x, y, 2.42, 1.75, CARD)
    shp = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.42), Inches(1.75))
    shp.fill.background(); shp.line.color.rgb = color; shp.line.width = Pt(1.2)
    rect(s, x, y, 2.42, 0.065, color)
    txt(s, name, x+0.1, y+0.1,  2.25, 0.52, size=11, bold=True, color=C_DARK)
    txt(s, val,  x+0.1, y+0.75, 2.25, 0.7,  size=13, bold=True, color=color)

rect(s, 0, 7.1, 13.33, 0.4, BG2)
txt(s, "Grafana  →  http://localhost:3001  |  admin / admin123  |  Dashboard : News Data Platform - Complet",
    0, 7.15, 13.33, 0.3, size=11, color=C_BLUE, align=PP_ALIGN.CENTER)

page_num(s, 13)


# ════════════════════════════════════════════════════════════
# SLIDE 14 — QUALITÉ DES DONNÉES
# ════════════════════════════════════════════════════════════
s = add_slide(); bg(s); left_bar(s, C_RED)
header(s, "Qualité des Données", "Contrôles automatisés — 3 dimensions", C_RED)

# Tableau des contrôles
rect(s, 0.4, 1.5, 9.0, 5.5, CARD)
shp = s.shapes.add_shape(1, Inches(0.4), Inches(1.5), Inches(9.0), Inches(5.5))
shp.fill.background(); shp.line.color.rgb = BORDER; shp.line.width = Pt(0.75)
txt(s, "Tests Implémentés", 0.6, 1.6, 8.7, 0.42, size=14, bold=True, color=C_RED)

for j, (h, w_) in enumerate([("Contrôle", 3.0), ("Dimension", 2.0), ("Règle", 3.8)]):
    x = 0.55 + sum([3.0, 2.0, 3.8][:j])
    rect(s, x-0.05, 2.1, w_+0.08, 0.42, BG2)
    txt(s, h, x, 2.15, w_, 0.35, size=11, bold=True, color=C_DARK)

checks = [
    ("Titre vide ou trop court", "Complétude",  C_RED,    "len(title) >= 5"),
    ("Date manquante",           "Complétude",  C_RED,    "published_at IS NOT NULL"),
    ("Contenu trop court",       "Validité",    C_ORANGE, "len(content) >= 100"),
    ("URL invalide",             "Validité",    C_ORANGE, "url.startswith('http')"),
    ("Source manquante",         "Complétude",  C_RED,    "source IS NOT NULL"),
    ("Langue incohérente",       "Cohérence",   C_PURPLE, "langue ↔ source attendue"),
]
for i, (ctrl, dim, color, rule) in enumerate(checks):
    y = 2.65 + i * 0.55
    txt(s, f"✗  {ctrl}", 0.55, y, 2.95, 0.5, size=11, color=C_TEXT)
    chip(s, dim, 3.55, y+0.06, 1.85, 0.36, color)
    txt(s, rule, 5.5,  y, 3.75, 0.5, size=11, color=C_MUTED, italic=True)

# Résultats
rect(s, 9.8, 1.5, 3.1, 5.5, CARD)
shp = s.shapes.add_shape(1, Inches(9.8), Inches(1.5), Inches(3.1), Inches(5.5))
shp.fill.background(); shp.line.color.rgb = BORDER; shp.line.width = Pt(0.75)
rect(s, 9.8, 1.5, 3.1, 0.07, C_GREEN)
txt(s, "Résultats", 9.95, 1.6, 2.85, 0.42, size=14, bold=True, color=C_GREEN)

res = [("Total", "173", C_DARK), ("Valides", "131", C_GREEN),
       ("Taux", "75.7%", C_ORANGE), ("Log", "PostgreSQL", C_BLUE)]
for i, (lbl, val, color) in enumerate(res):
    y = 2.2 + i * 1.15
    txt(s, lbl, 9.95, y,    2.85, 0.38, size=11, color=C_MUTED)
    txt(s, val, 9.95, y+0.42, 2.85, 0.58, size=18, bold=True, color=color)

page_num(s, 14)


# ════════════════════════════════════════════════════════════
# SLIDE 15 — GOUVERNANCE
# ════════════════════════════════════════════════════════════
s = add_slide(); bg(s); left_bar(s, C_PURPLE)
header(s, "Gouvernance des Données", "Traçabilité · Documentation · Transparence", C_PURPLE)

rect(s, 0.4, 1.5, 12.5, 1.1, BG2)
txt(s, "Data Lineage", 0.6, 1.58, 12.2, 0.35, size=12, bold=True, color=C_PURPLE)
lineage = "[EXTERNE] BBC·CNN·Reuters·Hespress·Akhbarona  →  [BRONZE] JSON  →  [SILVER] Parquet  →  [GOLD] Agrégats  →  [DWH] PostgreSQL"
txt(s, lineage, 0.6, 1.9, 12.2, 0.42, size=12, color=C_BLUE,
    italic=True, align=PP_ALIGN.CENTER)

pillars = [
    ("📋", "Documentation",    C_PURPLE,
     ["Schéma de chaque dataset", "Format + localisation", "Pipeline producteur", "governance/lineage.py"]),
    ("🔍", "Traçabilité",      C_BLUE,
     ["Source → destination", "lineage.jsonl log file", "Catalogue complet", "5 datasets documentés"]),
    ("📊", "Suivi Pipeline",   C_GREEN,
     ["Table pipeline_runs", "status + records_in/out", "Durée d'exécution", "Historique complet"]),
    ("✅", "Qualité",          C_ORANGE,
     ["data_quality_log", "Taux validation 75.7%", "Issues détectées", "Par pipeline + date"]),
]

for i, (icon, title, color, pts) in enumerate(pillars):
    x = 0.4 + i * 3.25
    rect(s, x, 2.8, 3.1, 4.3, CARD)
    shp = s.shapes.add_shape(1, Inches(x), Inches(2.8), Inches(3.1), Inches(4.3))
    shp.fill.background(); shp.line.color.rgb = color; shp.line.width = Pt(1.2)
    rect(s, x, 2.8, 3.1, 0.065, color)
    txt(s, icon,  x+0.15, 2.9,  0.65, 0.6, size=22)
    txt(s, title, x+0.85, 2.9,  2.1,  0.5, size=14, bold=True, color=color)
    for j, pt in enumerate(pts):
        txt(s, f"▶  {pt}", x+0.15, 3.55+j*0.72, 2.85, 0.65, size=12, color=C_TEXT)

page_num(s, 15)


# ════════════════════════════════════════════════════════════
# SLIDE 16 — DÉPLOIEMENT DOCKER
# ════════════════════════════════════════════════════════════
s = add_slide(); bg(s); left_bar(s, C_BLUE)
header(s, "Déploiement — Docker Compose", "11 services · 1 commande · Réseau isolé", C_BLUE)

services = [
    ("zookeeper",          C_TEAL,   "Coordination",    ":2181"),
    ("kafka",              C_TEAL,   "Message Broker",  ":9092"),
    ("kafka-ui",           C_TEAL,   "Interface",       ":8090"),
    ("minio",              C_ORANGE, "Object Storage",  ":9200"),
    ("postgres",           C_RED,    "Data Warehouse",  ":5432"),
    ("scraper",            C_GREEN,  "Web Scraping",    "interne"),
    ("consumer",           C_GREEN,  "Kafka→MinIO",     "interne"),
    ("airflow-webserver",  C_PURPLE, "Orchestration UI",":8080"),
    ("airflow-scheduler",  C_PURPLE, "DAG Scheduler",   "interne"),
    ("grafana",            C_TEAL,   "Dashboards",      ":3001"),
    ("prometheus",         C_ORANGE, "Métriques",       ":9090"),
]

for i, (name, color, desc, port) in enumerate(services):
    col = i % 4; row = i // 4
    x = 0.3 + col * 3.27; y = 1.58 + row * 1.7
    rect(s, x, y, 3.1, 1.52, CARD)
    shp = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(3.1), Inches(1.52))
    shp.fill.background(); shp.line.color.rgb = BORDER; shp.line.width = Pt(0.75)
    rect(s, x, y, 3.1, 0.065, color)
    txt(s, name,  x+0.12, y+0.1,  2.88, 0.42, size=13, bold=True, color=color)
    txt(s, desc,  x+0.12, y+0.6,  2.88, 0.38, size=11, color=C_TEXT)
    chip(s, port, x+0.12, y+1.08, 1.1,  0.33, BG2, C_BLUE)

rect(s, 0, 7.05, 13.33, 0.45, BG2)
txt(s, "docker compose up -d --build   ·   RAM : 5 GB minimum   ·   Docker Desktop 4.x",
    0, 7.12, 13.33, 0.32, size=12, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)

page_num(s, 16)


# ════════════════════════════════════════════════════════════
# SLIDE 17 — MONITORING
# ════════════════════════════════════════════════════════════
s = add_slide(); bg(s); left_bar(s, C_ORANGE)
header(s, "Monitoring — Prometheus + Grafana", "Observabilité complète de la plateforme", C_ORANGE)

for ci, (title, color, icon, pts, url) in enumerate([
    ("Prometheus", C_ORANGE, "🔭",
     ["Collecte métriques toutes les 15s",
      "CPU / RAM / réseau par container",
      "Métriques Kafka (messages/sec)",
      "Métriques PostgreSQL",
      "Alertes configurables",
      "Données temporelles (TSDB)"],
     "http://localhost:9090"),
    ("Grafana", C_TEAL, "📊",
     ["Source : PostgreSQL + Prometheus",
      "Dashboard News Data Platform",
      "13 panels analytiques live",
      "KPIs temps réel (auto-refresh)",
      "Provisioning automatique YAML",
      "Plugins : PostgreSQL natif"],
     "http://localhost:3001"),
]):
    x = 0.4 + ci * 6.5
    rect(s, x, 1.5, 6.2, 5.65, CARD)
    shp = s.shapes.add_shape(1, Inches(x), Inches(1.5), Inches(6.2), Inches(5.65))
    shp.fill.background(); shp.line.color.rgb = color; shp.line.width = Pt(1.5)
    rect(s, x, 1.5, 6.2, 0.07, color)
    txt(s, icon,  x+0.2,  1.6,  0.7,  0.6,  size=24)
    txt(s, title, x+0.95, 1.65, 5.0,  0.5,  size=20, bold=True, color=color)
    for i, pt in enumerate(pts):
        txt(s, f"✓  {pt}", x+0.2, 2.35+i*0.65, 5.8, 0.58, size=13, color=C_TEXT)
    rect(s, x+0.2, 6.45, 5.8, 0.48, BG2)
    txt(s, url,   x+0.3,  6.5,  5.6,  0.38, size=12, bold=True, color=color)

page_num(s, 17)


# ════════════════════════════════════════════════════════════
# SLIDE 18 — COMPARAISON TECHNOLOGIQUE
# ════════════════════════════════════════════════════════════
s = add_slide(); bg(s); left_bar(s, C_INDIGO)
header(s, "Choix Technologiques", "Pourquoi ces technologies ?", C_INDIGO)

comparisons = [
    ("Message Broker",   "Apache Kafka",  C_BLUE,   "RabbitMQ · Redis",   "Partitions · Rétention · Scalabilité"),
    ("Object Storage",   "MinIO",         C_ORANGE, "HDFS · AWS S3",      "API S3 · Local · Open-source · Docker"),
    ("Orchestration",    "Apache Airflow",C_GREEN,  "Apache NiFi · Cron", "DAGs Python · UI · Retry · Alertes"),
    ("Data Warehouse",   "PostgreSQL",    C_RED,    "MySQL · BigQuery",   "ACID · JSON · Extensions analytiques"),
    ("Visualisation",    "Grafana",       C_TEAL,   "PowerBI · Tableau",  "Open-source · PostgreSQL natif · Alertes"),
    ("Transformation",   "Python+Pandas", C_PURPLE, "Spark · dbt",        "Simple · Maîtrisé · Ecosystème riche"),
]

# En-têtes tableau
headers_t = ["Composant", "Choix", "Alternatives", "Justification"]
widths   = [2.3, 2.0, 2.5, 5.45]
x_starts = [0.4, 2.75, 4.8, 7.35]
for j, (h, x, w) in enumerate(zip(headers_t, x_starts, widths)):
    rect(s, x, 1.52, w-0.05, 0.48, C_INDIGO if j == 1 else BG2)
    txt(s, h, x+0.1, 1.57, w-0.15, 0.38, size=11, bold=True,
        color=C_WHITE if j == 1 else C_DARK)

for i, (comp, choice, color, alts, why) in enumerate(comparisons):
    y = 2.1 + i * 0.82
    bg_row = CARD if i % 2 == 0 else BG
    for x, w in zip(x_starts, widths):
        rect(s, x, y, w-0.05, 0.77, bg_row)
        shp = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w-0.05), Inches(0.77))
        shp.fill.background(); shp.line.color.rgb = BORDER; shp.line.width = Pt(0.4)
    txt(s, comp,   x_starts[0]+0.1, y+0.18, widths[0]-0.15, 0.45, size=12, color=C_TEXT)
    chip(s, choice, x_starts[1]+0.1, y+0.2, widths[1]-0.2, 0.38, color)
    txt(s, alts,   x_starts[2]+0.1, y+0.18, widths[2]-0.15, 0.45, size=11, color=C_MUTED)
    txt(s, why,    x_starts[3]+0.1, y+0.18, widths[3]-0.15, 0.45, size=11, color=C_TEXT)

page_num(s, 18)


# ════════════════════════════════════════════════════════════
# SLIDE 19 — DIFFICULTÉS ET SOLUTIONS
# ════════════════════════════════════════════════════════════
s = add_slide(); bg(s); left_bar(s, C_RED)
header(s, "Difficultés Rencontrées & Solutions", "Problèmes réels — approche méthodique", C_RED)

problems = [
    (C_RED,    "InconsistentClusterIdException Kafka",
               "Kafka refusait de démarrer après redémarrage Docker",
               "Suppression des volumes Docker corrompus + recréation du cluster"),
    (C_ORANGE, "Docker Desktop manquait de RAM",
               "Containers crashaient aléatoirement (OOM Killer)",
               "Augmentation limite RAM à 5 GB + limites mem_limit par service"),
    (C_PURPLE, "Grafana affichait 'No data'",
               "Datasource PostgreSQL mal configuré (uid conflit + jsonData.database manquant)",
               "Correction postgres.yml : suppression uid conflit + ajout jsonData.database"),
    (C_BLUE,   "Airflow réseau introuvable",
               "airflow-init ne voyait pas postgres / kafka au démarrage",
               "Ajout du réseau news-net dans airflow-common dans docker-compose.yml"),
]

for i, (color, title, prob, sol) in enumerate(problems):
    y = 1.55 + i * 1.42
    rect(s, 0.4, y, 12.5, 1.3, CARD)
    shp = s.shapes.add_shape(1, Inches(0.4), Inches(y), Inches(12.5), Inches(1.3))
    shp.fill.background(); shp.line.color.rgb = BORDER; shp.line.width = Pt(0.75)
    rect(s, 0.4, y, 0.07, 1.3, color)
    chip(s, f"#{i+1}", 0.55, y+0.1, 0.5, 0.38, color)
    txt(s, title, 1.15, y+0.08, 11.5, 0.38, size=14, bold=True, color=color)
    txt(s, f"❌  {prob}", 1.15, y+0.48, 11.5, 0.33, size=11, color=C_MUTED, italic=True)
    txt(s, f"✅  {sol}",  1.15, y+0.83, 11.5, 0.38, size=12, color=C_GREEN)

page_num(s, 19)


# ════════════════════════════════════════════════════════════
# SLIDE 20 — DÉMONSTRATION
# ════════════════════════════════════════════════════════════
s = add_slide(); bg(s); left_bar(s, C_GREEN)
header(s, "Démonstration Fonctionnelle", "Pipeline complet — bout en bout", C_GREEN)

steps_demo = [
    ("🕷", C_TEAL,   "Scraping", "Python\nBeautifulSoup"),
    ("📨", C_BLUE,   "Kafka",    "Streaming\nnews-articles"),
    ("🗄", C_ORANGE, "MinIO\nBronze", "900 fichiers\nJSON"),
    ("⚙️", C_INDIGO, "ETL\nSilver", "Parquet\nnettoyé"),
    ("🏛", C_RED,    "PostgreSQL\nDWH", "173 articles\nchargés"),
    ("📊", C_GREEN,  "Grafana\nDashboard", "13 panels\nLive"),
]

box_w = 1.75; gap_d = 0.28
total_w = len(steps_demo)*box_w + (len(steps_demo)-1)*gap_d
sx_d = (13.33 - total_w) / 2

for i, (icon, color, title, sub) in enumerate(steps_demo):
    x = sx_d + i * (box_w + gap_d)
    rect(s, x, 1.7, box_w, 2.0, CARD)
    shp = s.shapes.add_shape(1, Inches(x), Inches(1.7), Inches(box_w), Inches(2.0))
    shp.fill.background(); shp.line.color.rgb = color; shp.line.width = Pt(1.8)
    rect(s, x, 1.7, box_w, 0.07, color)
    txt(s, icon,  x+0.55, 1.82, 0.7, 0.6, size=22, align=PP_ALIGN.CENTER)
    txt(s, title, x+0.05, 2.47, box_w-0.1, 0.62, size=12, bold=True,
        color=color, align=PP_ALIGN.CENTER)
    txt(s, sub,   x+0.05, 3.1,  box_w-0.1, 0.5,  size=10,
        color=C_MUTED, align=PP_ALIGN.CENTER)
    if i < len(steps_demo)-1:
        txt(s, "→", x+box_w+0.04, 2.55, 0.22, 0.5, size=18, bold=True,
            color=C_MUTED, align=PP_ALIGN.CENTER)

interfaces = [("Airflow",  "localhost:8080", "admin/admin123",          C_PURPLE),
              ("MinIO",    "localhost:9200", "minioadmin/minioadmin123", C_ORANGE),
              ("Kafka UI", "localhost:8090", "—",                       C_TEAL),
              ("Grafana",  "localhost:3001", "admin/admin123",          C_GREEN),
              ("Metabase", "localhost:3000", "—",                       C_RED)]

txt(s, "Interfaces disponibles", 0.4, 3.95, 12.5, 0.38,
    size=12, bold=True, color=C_MUTED)

for i, (name, url, creds, color) in enumerate(interfaces):
    x = 0.35 + i * 2.59
    rect(s, x, 4.42, 2.45, 2.75, CARD)
    shp = s.shapes.add_shape(1, Inches(x), Inches(4.42), Inches(2.45), Inches(2.75))
    shp.fill.background(); shp.line.color.rgb = color; shp.line.width = Pt(1.2)
    rect(s, x, 4.42, 2.45, 0.065, color)
    txt(s, name,  x+0.12, 4.52, 2.25, 0.45, size=14, bold=True, color=color)
    txt(s, url,   x+0.12, 5.05, 2.25, 0.42, size=10, color=C_BLUE, italic=True)
    txt(s, creds, x+0.12, 5.52, 2.25, 0.42, size=10, color=C_MUTED)

rect(s, 0.35, 7.15, 12.65, 0.3, BG2)
txt(s, "Script de démonstration :   .\\demo.ps1   (PowerShell)",
    0.5, 7.18, 12.3, 0.25, size=11, color=C_BLUE, align=PP_ALIGN.CENTER)
page_num(s, 20)


# ════════════════════════════════════════════════════════════
# SLIDE 21 — CONCLUSION
# ════════════════════════════════════════════════════════════
s = add_slide(); bg(s); left_bar(s, C_BLUE)
header(s, "Conclusion", "Résumé · Résultats · Perspectives", C_BLUE)

rect(s, 0.4, 1.5, 4.0, 5.65, CARD)
shp = s.shapes.add_shape(1, Inches(0.4), Inches(1.5), Inches(4.0), Inches(5.65))
shp.fill.background(); shp.line.color.rgb = C_BLUE; shp.line.width = Pt(1.2)
rect(s, 0.4, 1.5, 4.0, 0.07, C_BLUE)
txt(s, "📊  Résultats", 0.6, 1.6, 3.7, 0.42, size=14, bold=True, color=C_BLUE)
for pt in ["900+ articles Bronze",
           "173 articles dans DWH",
           "5 sources actives",
           "131 mots-clés indexés",
           "3 pays couverts",
           "2 langues (EN / AR)",
           "75.7% taux de qualité",
           "13 panels Grafana live",
           "3 DAGs Airflow actifs",
           "11 services Docker"]:
    i = ["900+ articles Bronze","173 articles dans DWH","5 sources actives","131 mots-clés indexés",
         "3 pays couverts","2 langues (EN / AR)","75.7% taux de qualité","13 panels Grafana live",
         "3 DAGs Airflow actifs","11 services Docker"].index(pt)
    txt(s, f"✅  {pt}", 0.6, 2.18+i*0.51, 3.7, 0.48, size=12, color=C_TEXT)

rect(s, 4.75, 1.5, 3.9, 5.65, CARD)
shp = s.shapes.add_shape(1, Inches(4.75), Inches(1.5), Inches(3.9), Inches(5.65))
shp.fill.background(); shp.line.color.rgb = C_GREEN; shp.line.width = Pt(1.2)
rect(s, 4.75, 1.5, 3.9, 0.07, C_GREEN)
txt(s, "✅  Points Forts", 4.95, 1.6, 3.6, 0.42, size=14, bold=True, color=C_GREEN)
for i, pt in enumerate(["Architecture distribuée moderne",
                          "Pipeline 100% automatisé",
                          "Streaming temps réel (Kafka)",
                          "Qualité et gouvernance intégrées",
                          "Monitoring complet (Prometheus)",
                          "Déploiement reproductible Docker",
                          "Code versionné Git (55 fichiers)",
                          "Documentation complète",
                          "Kubernetes (optionnel) ✓",
                          "Bonus : Metabase BI"]):
    txt(s, f"▶  {pt}", 4.95, 2.18+i*0.51, 3.6, 0.48, size=12, color=C_TEXT)

rect(s, 9.0, 1.5, 3.9, 5.65, CARD)
shp = s.shapes.add_shape(1, Inches(9.0), Inches(1.5), Inches(3.9), Inches(5.65))
shp.fill.background(); shp.line.color.rgb = C_ORANGE; shp.line.width = Pt(1.2)
rect(s, 9.0, 1.5, 3.9, 0.07, C_ORANGE)
txt(s, "🚀  Perspectives", 9.2, 1.6, 3.6, 0.42, size=14, bold=True, color=C_ORANGE)
for i, pt in enumerate(["Déploiement Kubernetes",
                          "Ajout Al Jazeera + Barlamane",
                          "Analyse sentiment NLP",
                          "Détection fake news ML",
                          "Alertes tendances temps réel",
                          "API REST pour dashboards",
                          "Kafka multi-broker cluster",
                          "Tableau de bord mobile",
                          "Spark pour gros volumes",
                          "Intégration LLM (résumés)"]):
    txt(s, f"→  {pt}", 9.2, 2.18+i*0.51, 3.6, 0.48, size=12, color=C_TEXT)

page_num(s, 21)


# ════════════════════════════════════════════════════════════
# SLIDE 22 — LIVRABLES
# ════════════════════════════════════════════════════════════
s = add_slide(); bg(s); left_bar(s, C_GREEN)
header(s, "Livrables — Section 12", "Tous les livrables attendus remis ✅", C_GREEN)

livrables = [
    ("📁", "Code Source Git",         C_BLUE,   "55 fichiers · 2 commits\nmaster branch"),
    ("🏗",  "Schéma Architecture",    C_ORANGE, "docs/architecture.md\nDiagramme ASCII complet"),
    ("🐋", "Fichiers Docker",         C_TEAL,   "docker-compose.yml\n11 services définis"),
    ("☸️",  "Fichiers Kubernetes",    C_PURPLE, "k8s/ — 8 manifests YAML\nNamespace + Deployments"),
    ("📖", "Documentation Technique", C_RED,    "docs/technical_documentation.md\nInstallation + usage"),
    ("📊", "Dashboards Grafana",       C_GREEN,  "13 panels analytiques\nlocalhost:3001"),
    ("🎬", "Démonstration Pipeline",  C_INDIGO, "demo.ps1 PowerShell\nPipeline bout en bout"),
    ("📋", "Présentation",            C_ORANGE, "Ce fichier PPTX\n22 slides complets"),
]

for i, (icon, title, color, body) in enumerate(livrables):
    col = i % 4; row = i // 4
    x = 0.35 + col * 3.25; y = 1.58 + row * 2.65
    rect(s, x, y, 3.1, 2.45, CARD)
    shp = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(3.1), Inches(2.45))
    shp.fill.background(); shp.line.color.rgb = color; shp.line.width = Pt(1.5)
    rect(s, x, y, 3.1, 0.07, color)
    txt(s, icon,  x+0.18, y+0.15, 0.6,  0.6,  size=22)
    txt(s, "✓",  x+2.6,  y+0.18, 0.4,  0.42, size=16, bold=True,
        color=C_GREEN, align=PP_ALIGN.CENTER)
    txt(s, title, x+0.15, y+0.75, 2.9,  0.48, size=13, bold=True, color=color)
    txt(s, body,  x+0.15, y+1.3,  2.9,  1.0,  size=11, color=C_TEXT)

rect(s, 0, 7.1, 13.33, 0.4, BG2)
txt(s, "Projet Big Data complet · 11 services Docker · 173 articles analysés · Pipeline automatisé 24h/24",
    0, 7.15, 13.33, 0.3, size=12, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)

page_num(s, 22)


# ── SAUVEGARDE ──────────────────────────────────────────────
output = r"C:\Users\ayman\CascadeProjects\news-data-platform\docs\Presentation_News_Data_Platform.pptx"
prs.save(output)
print(f"Fichier sauvegarde : {output}")
print(f"Nombre de slides   : {len(prs.slides)}")
