"""
Intègre les vraies screenshots dans la présentation finale.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os, sys
sys.stdout.reconfigure(encoding='utf-8')

PPTX_IN  = r'C:\Users\ayman\CascadeProjects\news-data-platform\docs\Presentation_News_Data_Platform.pptx'
PPTX_OUT = r'C:\Users\ayman\CascadeProjects\news-data-platform\docs\Presentation_FINALE.pptx'
SS       = r'C:\Users\ayman\CascadeProjects\news-data-platform\docs\screenshots'

C_BLUE   = RGBColor(0x24, 0x63, 0xEB)
C_TEAL   = RGBColor(0x06, 0xB6, 0xD4)
C_GREEN  = RGBColor(0x10, 0xB9, 0x81)
C_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
C_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
C_DARK   = RGBColor(0x1E, 0x29, 0x4D)
C_MUTED  = RGBColor(0x7A, 0x8A, 0xAA)
BG       = RGBColor(0xF8, 0xF9, 0xFF)
BG2      = RGBColor(0xEE, 0xF2, 0xFF)

prs   = Presentation(PPTX_IN)
blank = prs.slide_layouts[6]


def add_screenshot_slide(title, subtitle, img_paths, accent, url, label=None):
    slide = prs.slides.add_slide(blank)

    bg_ = slide.background.fill; bg_.solid(); bg_.fore_color.rgb = BG

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()

    hdr = slide.shapes.add_shape(1, Inches(0.18), Inches(0), Inches(13.15), Inches(1.12))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = BG2; hdr.line.fill.background()

    sep = slide.shapes.add_shape(1, Inches(0.18), Inches(1.12), Inches(13.15), Inches(0.055))
    sep.fill.solid(); sep.fill.fore_color.rgb = accent; sep.line.fill.background()

    # Title
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.08), Inches(10.8), Inches(0.62))
    p  = tb.text_frame.paragraphs[0]; r = p.add_run()
    r.text = title; r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = C_DARK

    # Subtitle
    tb2 = slide.shapes.add_textbox(Inches(0.45), Inches(0.68), Inches(10.0), Inches(0.38))
    p2  = tb2.text_frame.paragraphs[0]; r2 = p2.add_run()
    r2.text = subtitle; r2.font.size = Pt(11); r2.font.italic = True; r2.font.color.rgb = C_MUTED

    # URL badge
    badge = slide.shapes.add_shape(1, Inches(11.0), Inches(0.22), Inches(2.1), Inches(0.58))
    badge.fill.solid(); badge.fill.fore_color.rgb = accent; badge.line.fill.background()
    tb3 = slide.shapes.add_textbox(Inches(11.02), Inches(0.26), Inches(2.06), Inches(0.48))
    p3  = tb3.text_frame.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER; r3 = p3.add_run()
    r3.text = url; r3.font.size = Pt(10); r3.font.bold = True
    r3.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Single image — full width
    if isinstance(img_paths, str):
        img_paths = [img_paths]

    if len(img_paths) == 1:
        p_ = img_paths[0]
        if os.path.exists(p_):
            slide.shapes.add_picture(p_, Inches(0.28), Inches(1.25), Inches(12.77), Inches(6.05))
            print(f"  ✓ {os.path.basename(p_)}")
    else:
        # Two images side by side
        for i, p_ in enumerate(img_paths):
            x = 0.28 + i * 6.45
            if os.path.exists(p_):
                slide.shapes.add_picture(p_, Inches(x), Inches(1.25), Inches(6.35), Inches(6.05))
                print(f"  ✓ {os.path.basename(p_)}")

    # Label tag
    if label:
        tag = slide.shapes.add_shape(1, Inches(0.28), Inches(7.0), Inches(3.0), Inches(0.35))
        tag.fill.solid(); tag.fill.fore_color.rgb = accent; tag.line.fill.background()
        tb4 = slide.shapes.add_textbox(Inches(0.32), Inches(7.02), Inches(2.9), Inches(0.3))
        p4  = tb4.text_frame.paragraphs[0]; r4 = p4.add_run()
        r4.text = label; r4.font.size = Pt(9); r4.font.bold = True
        r4.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


screenshots = [
    # (titre, sous-titre, image(s), couleur, url, label)
    ("📊  Grafana — Dashboard Analytique",
     "173 articles  ·  5 sources  ·  131 mots-clés  ·  75.7% qualité  ·  Dernier scraping : 14/05/2026 14:27",
     os.path.join(SS, "grafana_dashboard.png"),
     C_TEAL, "localhost:3001",
     "DÉMONSTRATION LIVE — Dashboard News Data Platform"),

    ("⚡  Airflow — 3 DAGs Actifs et Supervisés",
     "scraping_batch  ·  transformation_pipeline  ·  warehouse_load  ·  Active:3  Paused:0",
     os.path.join(SS, "airflow_dags.png"),
     C_GREEN, "localhost:8080",
     "ORCHESTRATION — Apache Airflow 2.8"),

    ("🗄  MinIO — Architecture Médaillon (Bronze / Gold / Silver)",
     "Bronze : 1012 objets · 4.0 MB  ·  Gold : 20 objets · 73.6 KB  ·  Silver : 14 objets · 4.2 MB",
     [os.path.join(SS, "minio_bronze.png"), os.path.join(SS, "minio_gold.png")],
     C_ORANGE, "localhost:9200",
     "DATA LAKE — MinIO Object Storage"),

    ("🗄  MinIO Silver — Données Nettoyées (Parquet)",
     "14 fichiers Parquet · 4.2 MB · dossier cleaned/ · Déduplication + Normalisation + Langue",
     os.path.join(SS, "minio_silver.png"),
     C_ORANGE, "localhost:9200",
     "DATA LAKE — Couche Silver"),

    ("📨  Kafka UI — Broker & Streaming",
     "1 Broker actif  ·  51 partitions Online  ·  5.78 MB  ·  51 segments  ·  Version 3.5-IV2",
     os.path.join(SS, "kafka_ui.png"),
     C_BLUE, "localhost:8090",
     "STREAMING — Apache Kafka"),
]

print("Ajout des slides screenshots (vraies images)...")
for args in screenshots:
    add_screenshot_slide(*args)

prs.save(PPTX_OUT)
print(f"\n✅ Fichier sauvegardé : {PPTX_OUT}")
print(f"   Total slides      : {len(prs.slides)}")
