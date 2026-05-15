from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os, sys
sys.stdout.reconfigure(encoding='utf-8')

PPTX_IN  = r'C:\Users\ayman\CascadeProjects\news-data-platform\docs\Presentation_News_Data_Platform.pptx'
PPTX_OUT = r'C:\Users\ayman\CascadeProjects\news-data-platform\docs\Presentation_FINAL_avec_Screenshots.pptx'
SS       = r'C:\Users\ayman\CascadeProjects\news-data-platform\docs\screenshots'

C_BLUE   = RGBColor(0x24, 0x63, 0xEB)
C_TEAL   = RGBColor(0x06, 0xB6, 0xD4)
C_GREEN  = RGBColor(0x10, 0xB9, 0x81)
C_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
C_DARK   = RGBColor(0x1E, 0x29, 0x4D)
C_MUTED  = RGBColor(0x7A, 0x8A, 0xAA)
BG       = RGBColor(0xF8, 0xF9, 0xFF)
BG2      = RGBColor(0xEE, 0xF2, 0xFF)

prs  = Presentation(PPTX_IN)
blank = prs.slide_layouts[6]

def add_screenshot_slide(title, subtitle, img_path, accent, url):
    slide = prs.slides.add_slide(blank)

    bg_ = slide.background.fill; bg_.solid(); bg_.fore_color.rgb = BG

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()

    hdr = slide.shapes.add_shape(1, Inches(0.18), Inches(0), Inches(13.15), Inches(1.15))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = BG2; hdr.line.fill.background()

    sep = slide.shapes.add_shape(1, Inches(0.18), Inches(1.15), Inches(13.15), Inches(0.055))
    sep.fill.solid(); sep.fill.fore_color.rgb = accent; sep.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.1), Inches(11.0), Inches(0.65))
    tf = tb.text_frame; p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = C_DARK

    tb2 = slide.shapes.add_textbox(Inches(0.45), Inches(0.72), Inches(9.0), Inches(0.38))
    tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]
    r2 = p2.add_run(); r2.text = subtitle
    r2.font.size = Pt(11); r2.font.italic = True; r2.font.color.rgb = C_MUTED

    url_box = slide.shapes.add_shape(1, Inches(11.1), Inches(0.25), Inches(2.05), Inches(0.55))
    url_box.fill.solid(); url_box.fill.fore_color.rgb = accent; url_box.line.fill.background()
    tb3 = slide.shapes.add_textbox(Inches(11.12), Inches(0.28), Inches(2.0), Inches(0.48))
    tf3 = tb3.text_frame; p3 = tf3.paragraphs[0]
    from pptx.enum.text import PP_ALIGN
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run(); r3.text = url
    r3.font.size = Pt(10); r3.font.bold = True
    r3.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    if os.path.exists(img_path):
        pic = slide.shapes.add_picture(img_path, Inches(0.28), Inches(1.28), Inches(12.77), Inches(6.0))
        print(f"  Image ajoutee : {os.path.basename(img_path)}")
    else:
        print(f"  ATTENTION image manquante : {img_path}")

screenshots = [
    ("📊  Grafana — Dashboard Analytique en Direct",
     "173 articles  ·  5 sources  ·  131 mots-clés  ·  13 panels  ·  données réelles",
     os.path.join(SS, "grafana_dashboard.png"), C_TEAL, "localhost:3001"),

    ("⚡  Airflow — Pipelines Orchestrés",
     "3 DAGs actifs  ·  scraping_batch @hourly  ·  transformation @daily  ·  warehouse @daily",
     os.path.join(SS, "airflow_home.png"), C_GREEN, "localhost:8080"),

    ("🗄  MinIO — Data Lake (Bronze / Silver / Gold)",
     "900+ fichiers JSON Bronze  ·  Parquet Silver  ·  Parquet Gold  ·  API S3 compatible",
     os.path.join(SS, "minio_buckets.png"), C_ORANGE, "localhost:9200"),

    ("📨  Kafka UI — Streaming Temps Réel",
     "Topic : news-articles  ·  3 partitions  ·  messages en continu depuis les scrapers",
     os.path.join(SS, "kafka_topic.png"), C_BLUE, "localhost:8090"),
]

print("Ajout des slides screenshots...")
for title, sub, img, accent, url in screenshots:
    add_screenshot_slide(title, sub, img, accent, url)

prs.save(PPTX_OUT)
print(f"\nFichier sauvegarde : {PPTX_OUT}")
print(f"Total slides       : {len(prs.slides)}")
